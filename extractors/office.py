"""OfficeExtractor — docx / xlsx / pptx (+ legacy doc/xls/ppt/odp).

Dependencies (pip packages):
    python-docx   → .docx (native, lightweight)
    openpyxl      → .xlsx (native, lightweight)
    python-pptx   → .pptx (native, lightweight)
    unstructured  → legacy .doc/.xls/.ppt + .odt/.ods/.odp fallback
                    (NOT installed by default here; requires libreoffice
                     binary on the container image; flagged at runtime)

Rationale: native libs are ~8 MiB total and bring no system deps. The
``unstructured`` backend pulls in ~400 MiB and needs a libreoffice
installation to convert the old binary formats, so we only import it
lazily if a legacy extension actually arrives.

Produces a multi-``Page`` output when the source has logical sections:

* pptx → one Page per slide (index = slide order, title = slide.title)
* xlsx → one Page per sheet (title = sheet name); rows are joined with
         tab separators. Truncated at ``max_rows_per_sheet`` config.
* docx → single-text output (no Pages); metadata still populated.
"""
from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Optional

from common.base_extractor import BaseExtractor
from common.models import ChangeEvent, ContentRef, ExtractedContent, Page


class OfficeExtractor(BaseExtractor):
    plugin_id = "office"
    version = "0.1.0"
    priority = 10

    supported_mime_types = [
        # docx
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        # xlsx
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        # pptx
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        # legacy binary office
        "application/msword",
        "application/vnd.ms-excel",
        "application/vnd.ms-powerpoint",
        # ODF
        "application/vnd.oasis.opendocument.text",
        "application/vnd.oasis.opendocument.spreadsheet",
        "application/vnd.oasis.opendocument.presentation",
    ]
    supported_extensions = [
        ".docx", ".doc",
        ".xlsx", ".xls",
        ".pptx", ".ppt",
        ".odt", ".ods", ".odp",
    ]

    _NATIVE_EXTS = {".docx", ".xlsx", ".pptx"}
    _LEGACY_EXTS = {".doc", ".xls", ".ppt", ".odt", ".ods", ".odp"}

    async def configure(self, config: dict) -> None:
        cfg = dict(config or {})
        self.max_rows_per_sheet: int = int(cfg.get("max_rows_per_sheet", 10000))
        self.include_speaker_notes: bool = bool(
            cfg.get("include_speaker_notes", True)
        )

    async def extract(
        self,
        content_ref: ContentRef,
        event: ChangeEvent,
    ) -> ExtractedContent:
        suffix = Path(event.path or "").suffix.lower()
        raw = await self._fetch_bytes(content_ref)

        if suffix == ".docx":
            return self._extract_docx(raw, event)
        if suffix == ".xlsx":
            return self._extract_xlsx(raw, event)
        if suffix == ".pptx":
            return self._extract_pptx(raw, event)
        if suffix in self._LEGACY_EXTS:
            return self._extract_legacy(raw, event, suffix)

        raise ValueError(
            f"OfficeExtractor cannot handle extension {suffix!r} for path "
            f"{event.path!r}"
        )

    # ── docx ────────────────────────────────────────────────────────

    def _extract_docx(self, raw: bytes, event: ChangeEvent) -> ExtractedContent:
        import docx  # python-docx

        doc = docx.Document(io.BytesIO(raw))
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        text = "\n".join(paragraphs)

        metadata = self._core_properties(doc.core_properties)
        metadata["format"] = "docx"

        return ExtractedContent(
            source_event_id=event.event_id,
            text=text,
            metadata=metadata,
            extractor_id=self.plugin_id,
            extractor_version=self.version,
        )

    # ── xlsx ────────────────────────────────────────────────────────

    def _extract_xlsx(self, raw: bytes, event: ChangeEvent) -> ExtractedContent:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        pages: list[Page] = []
        truncated: list[str] = []
        text_parts: list[str] = []

        for sheet_idx, sheet in enumerate(wb.worksheets):
            rows_text: list[str] = []
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= self.max_rows_per_sheet:
                    truncated.append(sheet.title)
                    break
                rows_text.append(
                    "\t".join("" if c is None else str(c) for c in row)
                )
            sheet_text = "\n".join(rows_text)
            pages.append(Page(index=sheet_idx, text=sheet_text, title=sheet.title))
            text_parts.append(f"# {sheet.title}\n{sheet_text}")

        metadata = self._core_properties(wb.properties)
        metadata["format"] = "xlsx"
        if truncated:
            metadata["truncated_sheets"] = truncated

        return ExtractedContent(
            source_event_id=event.event_id,
            text="\n\n".join(text_parts),
            pages=pages,
            metadata=metadata,
            extractor_id=self.plugin_id,
            extractor_version=self.version,
        )

    # ── pptx ────────────────────────────────────────────────────────

    def _extract_pptx(self, raw: bytes, event: ChangeEvent) -> ExtractedContent:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        pages: list[Page] = []
        text_parts: list[str] = []

        for slide_idx, slide in enumerate(prs.slides):
            lines: list[str] = []
            title: Optional[str] = None
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in para.runs)
                    if run_text:
                        lines.append(run_text)
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text or None
            if self.include_speaker_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    lines.append(f"[notes] {notes}")
            slide_text = "\n".join(lines)
            pages.append(Page(index=slide_idx, text=slide_text, title=title))
            text_parts.append(slide_text)

        metadata = self._core_properties(prs.core_properties)
        metadata["format"] = "pptx"

        return ExtractedContent(
            source_event_id=event.event_id,
            text="\n\n".join(text_parts),
            pages=pages,
            metadata=metadata,
            extractor_id=self.plugin_id,
            extractor_version=self.version,
        )

    # ── legacy / ODF (lazy-imported fallback) ───────────────────────

    def _extract_legacy(
        self,
        raw: bytes,
        event: ChangeEvent,
        suffix: str,
    ) -> ExtractedContent:
        try:
            from unstructured.partition.auto import partition  # type: ignore
        except Exception as exc:  # pragma: no cover - optional dep
            raise RuntimeError(
                f"OfficeExtractor: legacy extension {suffix!r} requires the "
                "optional 'unstructured' dependency (and libreoffice). "
                f"Import failed: {exc}"
            ) from exc

        elements = partition(file=io.BytesIO(raw))  # pragma: no cover
        text = "\n".join(str(el) for el in elements)  # pragma: no cover
        return ExtractedContent(  # pragma: no cover
            source_event_id=event.event_id,
            text=text,
            metadata={"format": suffix.lstrip("."), "backend": "unstructured"},
            extractor_id=self.plugin_id,
            extractor_version=self.version,
        )

    # ── helpers ─────────────────────────────────────────────────────

    @staticmethod
    def _core_properties(props: Any) -> dict[str, Any]:
        """Flatten openxml core_properties into a json-safe dict."""
        out: dict[str, Any] = {}
        for attr in ("author", "title", "subject", "last_modified_by"):
            val = getattr(props, attr, None)
            if val:
                out[attr] = str(val)
        created = getattr(props, "created", None)
        if created is not None:
            out["created_at"] = created.isoformat() if hasattr(created, "isoformat") else str(created)
        return out


__all__ = ["OfficeExtractor"]
